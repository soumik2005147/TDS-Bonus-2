
from typing import List, Dict, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

def _find_title_and_content_layout(prs: Presentation):
    # Try to find a layout with at least 2 placeholders (title + body)
    for layout in prs.slide_layouts:
        if len(layout.placeholders) >= 2:
            return layout
    # Fallback to first layout
    return prs.slide_layouts[0]

def _set_title(shape, text: str):
    if not shape or not getattr(shape, "text_frame", None):
        return
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    p.text = text.strip()[:200]

def _fill_bullets(shape, bullets: List[str]):
    if not shape or not getattr(shape, "text_frame", None):
        return
    tf = shape.text_frame
    tf.clear()
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = b.strip()
        p.level = 0

def _add_notes(slide, notes: Optional[str]):
    if not notes:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    tf.text = notes.strip()[:2000]

def build_ppt(template_path: str, slides: List[Dict], output_path: str):
    prs = Presentation(template_path)
    layout = _find_title_and_content_layout(prs)

    for s in slides:
        title = s.get("title", "").strip() or " "
        bullets = s.get("bullets", []) or []
        notes = s.get("notes", None)

        slide = prs.slides.add_slide(layout)

        # Assume placeholder 0 is title, 1 is body (common in many templates)
        title_shape = slide.shapes.title
        body_shape = None
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]

        _set_title(title_shape, title)
        _fill_bullets(body_shape, bullets[:6])
        _add_notes(slide, notes)

    prs.save(output_path)
